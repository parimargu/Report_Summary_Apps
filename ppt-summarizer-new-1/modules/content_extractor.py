"""
Content extraction module.

Extracts text, tables, and other content from PowerPoint slides.
"""

from typing import Dict, List, Optional, Tuple
from dataclasses import dataclass

import pandas as pd
from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from pptx.table import Table
from pptx.text.text import TextFrame

from modules.logger import get_logger
from modules.config_manager import get_config

logger = get_logger(__name__)
config = get_config()


@dataclass
class SlideContent:
    """Data class to store extracted slide content."""
    slide_number: int
    title: Optional[str]
    text_content: List[str]
    tables: List[pd.DataFrame]
    table_texts: List[str]  # Raw table text for LLM
    has_content: bool


class ContentExtractor:
    """Extracts content from PowerPoint presentations."""

    def __init__(self):
        """Initialize content extractor with configuration."""
        self.min_table_rows = config.get("extraction.min_table_rows", 2)
        self.min_table_cols = config.get("extraction.min_table_cols", 2)
        self.preserve_formatting = config.get("extraction.preserve_formatting", True)
        logger.info("ContentExtractor initialized")

    def extract_all_slides(self, presentation: Presentation) -> List[SlideContent]:
        """
        Extract content from all slides in presentation.

        Args:
            presentation: PowerPoint presentation object

        Returns:
            List of SlideContent objects
        """
        all_slides = []

        try:
            for idx, slide in enumerate(presentation.slides, start=1):
                logger.debug(f"Extracting content from slide {idx}")
                slide_content = self.extract_slide_content(slide, idx)
                all_slides.append(slide_content)

            logger.info(f"Extracted content from {len(all_slides)} slides")
            return all_slides

        except Exception as e:
            logger.error(f"Error extracting slides: {str(e)}", exc_info=True)
            return []

    def extract_slide_content(self, slide, slide_number: int) -> SlideContent:
        """
        Extract content from a single slide.

        Args:
            slide: PowerPoint slide object
            slide_number: Slide number (1-indexed)

        Returns:
            SlideContent object with extracted data
        """
        try:
            # Extract title
            title = self._extract_title(slide)

            # Extract text content
            text_content = self._extract_text(slide)

            # Extract tables
            tables, table_texts = self._extract_tables(slide)

            # Check if slide has any content
            has_content = bool(title or text_content or tables)

            slide_content = SlideContent(
                slide_number=slide_number,
                title=title,
                text_content=text_content,
                tables=tables,
                table_texts=table_texts,
                has_content=has_content
            )

            logger.debug(
                f"Slide {slide_number}: title={bool(title)}, "
                f"text_blocks={len(text_content)}, tables={len(tables)}"
            )

            return slide_content

        except Exception as e:
            logger.error(f"Error extracting content from slide {slide_number}: {str(e)}")
            return SlideContent(
                slide_number=slide_number,
                title=None,
                text_content=[],
                tables=[],
                table_texts=[],
                has_content=False
            )

    def _extract_title(self, slide) -> Optional[str]:
        """
        Extract slide title.

        Args:
            slide: PowerPoint slide object

        Returns:
            Slide title or None
        """
        try:
            if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
                if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                    title = slide.shapes.title.text.strip()
                    if title:
                        return title
        except Exception as e:
            logger.debug(f"Could not extract title: {str(e)}")

        return None

    def _extract_text(self, slide) -> List[str]:
        """
        Extract text content from slide (excluding tables and title).

        Args:
            slide: PowerPoint slide object

        Returns:
            List of text strings
        """
        text_blocks = []

        try:
            for shape in slide.shapes:
                # Skip title shape
                if hasattr(slide.shapes, 'title') and shape == slide.shapes.title:
                    continue

                # Skip table shapes
                if shape.has_table:
                    continue

                # Extract text from text frames
                if hasattr(shape, 'text_frame'):
                    text = self._extract_text_from_frame(shape.text_frame)
                    if text:
                        text_blocks.append(text)

                # Extract text from text property
                elif hasattr(shape, 'text'):
                    text = shape.text.strip()
                    if text:
                        text_blocks.append(text)

        except Exception as e:
            logger.error(f"Error extracting text: {str(e)}")

        return text_blocks

    def _extract_text_from_frame(self, text_frame: TextFrame) -> str:
        """
        Extract text from a text frame.

        Args:
            text_frame: PowerPoint text frame object

        Returns:
            Extracted text
        """
        try:
            paragraphs = []
            for paragraph in text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    paragraphs.append(para_text)

            return "\n".join(paragraphs)
        except Exception as e:
            logger.debug(f"Error extracting text from frame: {str(e)}")
            return ""

    def _extract_tables(self, slide) -> Tuple[List[pd.DataFrame], List[str]]:
        """
        Extract tables from slide.

        Args:
            slide: PowerPoint slide object

        Returns:
            Tuple of (list of DataFrames, list of formatted table strings)
        """
        tables = []
        table_texts = []

        try:
            for shape in slide.shapes:
                if shape.has_table:
                    df, table_text = self._convert_table_to_dataframe(shape.table)

                    # Only include tables that meet minimum size requirements
                    if df is not None and not df.empty:
                        if len(df) >= self.min_table_rows and len(df.columns) >= self.min_table_cols:
                            tables.append(df)
                            table_texts.append(table_text)
                            logger.debug(f"Extracted table with shape: {df.shape}")

        except Exception as e:
            logger.error(f"Error extracting tables: {str(e)}")

        return tables, table_texts

    def _convert_table_to_dataframe(self, table: Table) -> Tuple[Optional[pd.DataFrame], str]:
        """
        Convert PowerPoint table to pandas DataFrame.

        Args:
            table: PowerPoint table object

        Returns:
            Tuple of (DataFrame, formatted table string for LLM)
        """
        try:
            # Extract all cell values
            data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    row_data.append(cell_text)
                data.append(row_data)

            if not data:
                return None, ""

            # Create DataFrame
            # Try to use first row as headers if it looks like headers
            if len(data) > 1:
                headers = data[0]
                # Check if first row looks like headers (non-numeric or descriptive)
                if any(not self._is_numeric(cell) for cell in headers):
                    df = pd.DataFrame(data[1:], columns=headers)
                else:
                    df = pd.DataFrame(data)
            else:
                df = pd.DataFrame(data)

            # Create formatted string for LLM
            table_text = self._format_table_for_llm(df)

            return df, table_text

        except Exception as e:
            logger.error(f"Error converting table to DataFrame: {str(e)}")
            return None, ""

    def _format_table_for_llm(self, df: pd.DataFrame) -> str:
        """
        Format DataFrame as readable text for LLM processing.

        Args:
            df: pandas DataFrame

        Returns:
            Formatted table string
        """
        try:
            # Use pandas to_markdown for nice formatting
            # Fallback to to_string if markdown is not available
            try:
                table_str = df.to_markdown(index=False)
            except AttributeError:
                table_str = df.to_string(index=False)

            return table_str

        except Exception as e:
            logger.error(f"Error formatting table: {str(e)}")
            return str(df)

    def _is_numeric(self, value: str) -> bool:
        """
        Check if a string value is numeric.

        Args:
            value: String to check

        Returns:
            True if numeric, False otherwise
        """
        try:
            # Remove common formatting characters
            cleaned = value.replace(',', '').replace('$', '').replace('%', '').strip()
            float(cleaned)
            return True
        except (ValueError, AttributeError):
            return False

    def get_slide_summary(self, slide_content: SlideContent) -> Dict[str, any]:
        """
        Get summary statistics for a slide.

        Args:
            slide_content: SlideContent object

        Returns:
            Dictionary with summary information
        """
        return {
            "slide_number": slide_content.slide_number,
            "has_title": slide_content.title is not None,
            "text_blocks": len(slide_content.text_content),
            "table_count": len(slide_content.tables),
            "has_content": slide_content.has_content,
            "total_text_length": sum(len(text) for text in slide_content.text_content)
        }


# Example usage
if __name__ == "__main__":
    from modules.file_parser import FileParser

    parser = FileParser()
    extractor = ContentExtractor()

    try:
        prs = parser.parse_presentation("sample_data/sample_loan_forecast.pptx")
        slides = extractor.extract_all_slides(prs)

        for slide in slides:
            summary = extractor.get_slide_summary(slide)
            print(f"\nSlide {summary['slide_number']}:")
            print(f"  Title: {slide.title}")
            print(f"  Text blocks: {summary['text_blocks']}")
            print(f"  Tables: {summary['table_count']}")

    except Exception as e:
        print(f"Error: {e}")
