import pytest
from modules import ppt_parser
import io
from pptx import Presentation

def create_sample_ppt():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_textbox(0,0,100,100).text = "Test paragraph"
    table = slide.shapes.add_table(2,2,0,0,100,100).table
    table.cell(0,0).text = "Segment"
    table.cell(0,1).text = "Loan Default Rate (%)"
    table.cell(1,0).text = "Retail"
    table.cell(1,1).text = "-2"
    f = io.BytesIO()
    prs.save(f)
    f.seek(0)
    return f

def test_parse_ppt():
    ppt_file = create_sample_ppt()
    slides_content = ppt_parser.parse_ppt(ppt_file)
    assert len(slides_content) == 1
    assert slides_content[0]['paragraphs'][0] == "Test paragraph"
    assert slides_content[0]['tables'][0].iloc[0,1] == "-2"
