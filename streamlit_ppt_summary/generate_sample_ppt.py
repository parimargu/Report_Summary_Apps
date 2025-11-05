from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Slide 1 - Introduction
slide1 = prs.slides.add_slide(prs.slide_layouts[5])
tx_box1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
tf1 = tx_box1.text_frame
p1 = tf1.add_paragraph()
p1.text = "Quarterly Loan Forecast Report - Commercial Banking Segment"
p1.font.size = Pt(24)
p1.font.bold = True

# Slide 2 - Table with loan metrics
slide2 = prs.slides.add_slide(prs.slide_layouts[5])
tx_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
tf2 = tx_box2.text_frame
tf2.text = "Loan Performance Metrics (Q3 2025)"

# Add table
rows, cols = 4, 3
table_shape = slide2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(3))
table = table_shape.table

# Table headers
table.cell(0,0).text = "Segment"
table.cell(0,1).text = "Loan Default Rate (%)"
table.cell(0,2).text = "Net Rate (%)"

# Sample data
data = [
    ["Retail", "-2", "5"],
    ["Corporate", "7", "-1"],
    ["SME", "3", "2"]
]

# Fill table
for r, row in enumerate(data, start=1):
    for c, val in enumerate(row):
        table.cell(r, c).text = val
        # Highlight negative values in red and bold
        if c > 0 and float(val) < 0:
            cell_text = table.cell(r, c).text_frame.paragraphs[0].runs[0]
            cell_text.font.color.rgb = RGBColor(255, 0, 0)
            cell_text.font.bold = True

# Slide 3 - Summary Text
slide3 = prs.slides.add_slide(prs.slide_layouts[5])
tx_box3 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
tx_box3.text_frame.text = "This presentation contains forecasts and risk indicators for various banking segments."

# Save the presentation
prs.save("sample_loan_forecast.pptx")
print("Sample PPTX created: sample_loan_forecast.pptx")
