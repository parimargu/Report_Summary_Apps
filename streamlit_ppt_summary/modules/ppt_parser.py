from pptx import Presentation
import pandas as pd
import logging

logger = logging.getLogger(__name__)


def parse_ppt(file):
    prs = Presentation(file)
    slides_content = []

    for slide in prs.slides:
        slide_dict = {'paragraphs': [], 'tables': []}

        # Extract paragraphs
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_dict['paragraphs'].append(shape.text)

        # Extract tables
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                data = [[cell.text for cell in row.cells] for row in table.rows]
                df = pd.DataFrame(data[1:], columns=data[0])
                slide_dict['tables'].append(df)

        slides_content.append(slide_dict)
    logger.info("Parsed %d slides", len(slides_content))
    return slides_content
