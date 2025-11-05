"""
PowerPoint file parsing module.

Handles opening and reading PowerPoint presentations (.ppt and .pptx files).
"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType

from modules.logger import get_logger
from modules.config_manager import get_config

logger = get_logger(__name__)
config = get_config()


class FileParser:
    """Handles PowerPoint file parsing and validation."""

    def __init__(self):
        """Initialize file parser with configuration."""
        self.max_file_size_mb = config.get("app.max_file_size_mb", 5)
        self.supported_formats = config.get("app.supported_formats", [".ppt", ".pptx"])
        logger.info(f"FileParser initialized with max size: {self.max_file_size_mb}MB")

    def validate_file(self, file_path: str) -> tuple[bool, Optional[str]]:
        """
        Validate uploaded file.

        Args:
            file_path: Path to the file to validate

        Returns:
            Tuple of (is_valid, error_message)
            - (True, None) if file is valid
            - (False, error_message) if validation fails
        """
        try:
            path = Path(file_path)

            # Check if file exists
            if not path.exists():
                error_msg = f"File not found: {file_path}"
                logger.error(error_msg)
                return False, error_msg

            # Check file extension
            if path.suffix.lower() not in self.supported_formats:
                error_msg = f"Unsupported file format: {path.suffix}. Supported: {', '.join(self.supported_formats)}"
                logger.error(error_msg)
                return False, error_msg

            # Check file size
            file_size_mb = path.stat().st_size / (1024 * 1024)
            if file_size_mb > self.max_file_size_mb:
                error_msg = f"File size ({file_size_mb:.2f}MB) exceeds maximum allowed size ({self.max_file_size_mb}MB)"
                logger.error(error_msg)
                return False, error_msg

            logger.info(f"File validation passed: {file_path} ({file_size_mb:.2f}MB)")
            return True, None

        except Exception as e:
            error_msg = f"Error validating file: {str(e)}"
            logger.error(error_msg, exc_info=True)
            return False, error_msg

    def validate_uploaded_file(self, uploaded_file) -> tuple[bool, Optional[str]]:
        """
        Validate uploaded file from Streamlit.

        Args:
            uploaded_file: Streamlit UploadedFile object

        Returns:
            Tuple of (is_valid, error_message)
        """
        try:
            # Check file name extension
            file_name = uploaded_file.name
            file_extension = Path(file_name).suffix.lower()

            if file_extension not in self.supported_formats:
                error_msg = f"Unsupported file format: {file_extension}. Supported: {', '.join(self.supported_formats)}"
                logger.error(error_msg)
                return False, error_msg

            # Check file size
            file_size_mb = uploaded_file.size / (1024 * 1024)
            if file_size_mb > self.max_file_size_mb:
                error_msg = f"File size ({file_size_mb:.2f}MB) exceeds maximum allowed size ({self.max_file_size_mb}MB)"
                logger.error(error_msg)
                return False, error_msg

            logger.info(f"Uploaded file validation passed: {file_name} ({file_size_mb:.2f}MB)")
            return True, None

        except Exception as e:
            error_msg = f"Error validating uploaded file: {str(e)}"
            logger.error(error_msg, exc_info=True)
            return False, error_msg

    def parse_presentation(self, file_path: str) -> Optional[PresentationType]:
        """
        Parse PowerPoint presentation file.

        Args:
            file_path: Path to PowerPoint file

        Returns:
            Presentation object if successful, None if parsing fails

        Raises:
            FileNotFoundError: If file doesn't exist
            ValueError: If file format is invalid
        """
        try:
            # Validate file first
            is_valid, error_msg = self.validate_file(file_path)
            if not is_valid:
                raise ValueError(error_msg)

            logger.info(f"Parsing presentation: {file_path}")
            presentation = Presentation(file_path)

            slide_count = len(presentation.slides)
            logger.info(f"Successfully parsed presentation with {slide_count} slides")

            return presentation

        except FileNotFoundError as e:
            logger.error(f"File not found: {file_path}")
            raise

        except Exception as e:
            logger.error(f"Error parsing presentation: {str(e)}", exc_info=True)
            raise ValueError(f"Failed to parse PowerPoint file: {str(e)}")

    def parse_uploaded_presentation(self, uploaded_file) -> Optional[PresentationType]:
        """
        Parse PowerPoint presentation from Streamlit uploaded file.

        Args:
            uploaded_file: Streamlit UploadedFile object

        Returns:
            Presentation object if successful, None if parsing fails
        """
        try:
            # Validate uploaded file
            is_valid, error_msg = self.validate_uploaded_file(uploaded_file)
            if not is_valid:
                raise ValueError(error_msg)

            logger.info(f"Parsing uploaded presentation: {uploaded_file.name}")

            # Read file content and parse
            presentation = Presentation(uploaded_file)

            slide_count = len(presentation.slides)
            logger.info(f"Successfully parsed uploaded presentation with {slide_count} slides")

            return presentation

        except Exception as e:
            logger.error(f"Error parsing uploaded presentation: {str(e)}", exc_info=True)
            raise ValueError(f"Failed to parse PowerPoint file: {str(e)}")

    def get_slide_count(self, presentation: PresentationType) -> int:
        """
        Get number of slides in presentation.

        Args:
            presentation: Presentation object

        Returns:
            Number of slides
        """
        try:
            count = len(presentation.slides)
            logger.debug(f"Presentation contains {count} slides")
            return count
        except Exception as e:
            logger.error(f"Error getting slide count: {str(e)}")
            return 0


# Example usage
if __name__ == "__main__":
    parser = FileParser()

    # Test validation
    is_valid, error = parser.validate_file("sample_data/sample_loan_forecast.pptx")
    print(f"Validation result: {is_valid}, Error: {error}")

    # Test parsing
    if is_valid:
        try:
            prs = parser.parse_presentation("sample_data/sample_loan_forecast.pptx")
            print(f"Slides: {parser.get_slide_count(prs)}")
        except Exception as e:
            print(f"Parsing error: {e}")
